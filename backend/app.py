from flask import Flask, request, jsonify
from flask_cors import CORS
from markitdown import MarkItDown
import os
import tempfile
from werkzeug.utils import secure_filename
from PIL import Image
import pdfplumber
from docx import Document
import openpyxl
from pptx import Presentation

# Try to import Tesseract and configure path
try:
    import pytesseract
    import subprocess
    
    # Common Tesseract installation paths on Windows
    possible_paths = [
        r'C:\Program Files\Tesseract-OCR\tesseract.exe',
        r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
        r'C:\Users\%USERNAME%\AppData\Local\Programs\Tesseract-OCR\tesseract.exe',
        'tesseract'  # If it's in PATH
    ]
    
    TESSERACT_AVAILABLE = False
    
    for path in possible_paths:
        try:
            # Expand environment variables
            expanded_path = os.path.expandvars(path)
            
            if path == 'tesseract':
                # Test if tesseract is in PATH
                result = subprocess.run(['tesseract', '--version'], 
                                      capture_output=True, text=True, timeout=5)
                if result.returncode == 0:
                    TESSERACT_AVAILABLE = True
                    print(f"Tesseract found in PATH")
                    break
            else:
                # Test specific path
                if os.path.exists(expanded_path):
                    pytesseract.pytesseract.tesseract_cmd = expanded_path
                    # Test it works
                    result = subprocess.run([expanded_path, '--version'], 
                                          capture_output=True, text=True, timeout=5)
                    if result.returncode == 0:
                        TESSERACT_AVAILABLE = True
                        print(f"Tesseract found at: {expanded_path}")
                        break
        except Exception as e:
            continue
    
    if not TESSERACT_AVAILABLE:
        print("Tesseract installed but not working - check installation")
        
except ImportError:
    TESSERACT_AVAILABLE = False
    print("pytesseract not installed")

app = Flask(__name__)
CORS(app)

UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'txt', 'pdf', 'png', 'jpg', 'jpeg', 'gif', 'docx', 'pptx', 'xlsx'}

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def clean_vietnamese_text(text):
    """Clean and normalize Vietnamese OCR text"""
    import re
    
    # Remove excessive whitespace
    text = re.sub(r'\s+', ' ', text)
    
    # Common OCR corrections for Vietnamese
    corrections = {
        'ii': 'ì',
        'ifi': 'ỉ', 
        'ỉỉ': 'ủ',
        'rii': 'ủ',
        'c&n': 'căn',
        'C&n': 'Căn',
        'ctr': 'cứ',
        'Ctr': 'Cứ',
        'ném': 'năm',
        'N6i': 'Nội',
        'n6i': 'nội',
        'dé': 'đề',
        'Dé': 'Đề',
        'nghj': 'nghị',
        'Nghj': 'Nghị',
        'ciia': 'của',
        'Ciia': 'Của',
        'B6': 'Bộ',
        'b6': 'bộ',
        'truvéng': 'trưởng',
        'Truvéng': 'Trưởng',
        'héa': 'hóa',
        'Héa': 'Hóa',
        "vu'c": 'vực',
        "Vu'c": 'Vực',
        'inh': 'lĩnh',
        'Inh': 'Lĩnh',
        "vy'c": 'vực',
        "Vy'c": 'Vực',
        'ly': 'lý',
        'Ly': 'Lý',
        'phat': 'phạt',
        'Phat': 'Phạt',
        'hanh': 'hành',
        'Hanh': 'Hành',
        'chinh': 'chính',
        'Chinh': 'Chính',
        'dinh': 'định',
        'Dinh': 'Định',
        'quy': 'quy',
        'Quy': 'Quy',
        'xir': 'xử',
        'Xir': 'Xử',
        'pham': 'phạm',
        'Pham': 'Phạm',
        'lich': 'lịch',
        'Lich': 'Lịch',
        'luét': 'luật',
        'Luét': 'Luật',
        'chee': 'chức',
        'Chee': 'Chức',
        'phi': 'phủ',
        'Phi': 'Phủ',
        'lap': 'lập',
        'Lap': 'Lập',
        'oc': 'độc',
        'Oc': 'Độc',
        'ty': 'tự',
        'Ty': 'Tự',
        'phic': 'phúc',
        'Phic': 'Phúc',
        'phủc': 'phúc',
        'Phủc': 'Phúc',
        't6': 'tổ',
        'T6': 'Tổ',
        'Chlĩnh': 'Chính',
        'chlĩnh': 'chính',
        'dlĩnh': 'định',
        'Dlĩnh': 'Định',
        'vue': 'vực',
        'Vue': 'Vực',
        'Luat': 'Luật',
        'luat': 'luật',
        'cìa': 'của',
        'Cìa': 'Của',
        'Van': 'Văn',
        'van': 'văn',
        'Thé': 'Thể',
        'thé': 'thể',
        'va': 'và',
        'Va': 'Và',
        'ban': 'ban',
        'Ban': 'Ban',
        'ngay': 'ngày',
        'Ngay': 'Ngày',
        'thang': 'tháng',
        'Thang': 'Tháng',
        'nam': 'năm',
        'Nam': 'Năm',
        'xu': 'xử',
        'Xu': 'Xử',
        'XU': 'XỬ'
    }
    
    for wrong, correct in corrections.items():
        text = text.replace(wrong, correct)
    
    return text.strip()

def extract_text_fallback(file_path, file_extension):
    """Fallback text extraction methods if MarkItDown fails"""
    print(f"Trying fallback for file: {file_path}, extension: {file_extension}")
    try:
        if file_extension.lower() in ['.png', '.jpg', '.jpeg', '.gif']:
            # OCR for images
            if not TESSERACT_AVAILABLE:
                return "Tesseract OCR not installed. See INSTALL_OCR.md for setup instructions."
            
            try:
                image = Image.open(file_path)
                
                # Try Vietnamese first, then English, then both
                languages_to_try = ['vie', 'eng', 'vie+eng']
                
                for lang in languages_to_try:
                    try:
                        # Configure OCR with better settings for Vietnamese
                        custom_config = r'--oem 3 --psm 6'
                        
                        text = pytesseract.image_to_string(
                            image, 
                            lang=lang,
                            config=custom_config
                        )
                        
                        if text.strip():
                            print(f"OCR successful with language: {lang}")
                            return text.strip()
                    except Exception as lang_error:
                        print(f"OCR failed with {lang}: {lang_error}")
                        continue
                
                # If all language attempts failed, try basic OCR
                text = pytesseract.image_to_string(image)
                return text.strip() if text.strip() else "No text found in image."
                
            except Exception as e:
                return f"OCR failed: {str(e)}"
        
        elif file_extension.lower() == '.pdf':
            # PDF text extraction
            print("Attempting PDF text extraction...")
            text = ""
            try:
                with pdfplumber.open(file_path) as pdf:
                    print(f"PDF has {len(pdf.pages)} pages")
                    for page_num, page in enumerate(pdf.pages):
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                            print(f"Page {page_num + 1} extracted {len(page_text)} characters")
                
                print(f"Total PDF text length: {len(text)}")
                
                if not text.strip():
                    if TESSERACT_AVAILABLE:
                        return "This PDF contains scanned content. OCR for PDF scans is not implemented yet. Try converting PDF pages to images first."
                    else:
                        return "This PDF contains scanned content. Install Tesseract OCR (see INSTALL_OCR.md) to extract text from scanned PDFs."
                
                return text
            except Exception as e:
                print(f"PDF extraction error: {e}")
                return f"Error reading PDF: {str(e)}"
        
        elif file_extension.lower() == '.docx':
            # Word document
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        
        elif file_extension.lower() == '.xlsx':
            # Excel file
            workbook = openpyxl.load_workbook(file_path)
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
                text += "\n"
            return text
        
        elif file_extension.lower() == '.pptx':
            # PowerPoint
            prs = Presentation(file_path)
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + "\n"
                text += "\n"
            return text
        
        elif file_extension.lower() == '.txt':
            # Plain text
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
    
    except Exception as e:
        print(f"Fallback extraction failed: {e}")
        return None
    
    return None

@app.route('/api/health', methods=['GET'])
def health_check():
    return jsonify({'status': 'OK', 'message': 'OCR API is running'})

@app.route('/api/convert', methods=['POST'])
def convert_document():
    temp_file_path = None
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            
            # Create a temporary file
            temp_fd, temp_file_path = tempfile.mkstemp(suffix=os.path.splitext(filename)[1])
            try:
                # Write file data to temporary file
                with os.fdopen(temp_fd, 'wb') as temp_file:
                    file.save(temp_file)
                
                # Initialize MarkItDown
                md = MarkItDown()
                
                # Convert the document
                result = md.convert(temp_file_path)
                
                # Get extracted text
                extracted_text = result.text_content if result.text_content else ""
                
                # Debug logging
                print(f"MarkItDown result length: {len(extracted_text)}")
                
                # If MarkItDown failed, try fallback methods
                if not extracted_text.strip():
                    print("MarkItDown failed, trying fallback methods...")
                    extracted_text = extract_text_fallback(temp_file_path, os.path.splitext(filename)[1])
                    if extracted_text:
                        print(f"Fallback extraction successful, length: {len(extracted_text)}")
                    else:
                        print("Fallback extraction also failed")
                        extracted_text = "No text could be extracted from this file."
                
                return jsonify({
                    'success': True,
                    'text': extracted_text,
                    'filename': filename
                })
            finally:
                # Clean up temporary file
                if temp_file_path and os.path.exists(temp_file_path):
                    try:
                        os.unlink(temp_file_path)
                    except:
                        pass  # Ignore cleanup errors
        
        return jsonify({'error': 'File type not allowed'}), 400
        
    except Exception as e:
        # Clean up temporary file in case of error
        if temp_file_path and os.path.exists(temp_file_path):
            try:
                os.unlink(temp_file_path)
            except:
                pass
        return jsonify({'error': str(e)}), 500

@app.route('/api/supported-formats', methods=['GET'])
def supported_formats():
    return jsonify({
        'formats': list(ALLOWED_EXTENSIONS),
        'description': 'Supported file formats for OCR conversion'
    })

@app.route('/api/test-ocr', methods=['GET'])
def test_ocr():
    """Test OCR functionality"""
    try:
        # Test Tesseract
        import subprocess
        result = subprocess.run(['tesseract', '--version'], capture_output=True, text=True)
        tesseract_version = result.stdout.split('\n')[0] if result.returncode == 0 else "Not found"
        
        return jsonify({
            'tesseract_available': result.returncode == 0,
            'tesseract_version': tesseract_version,
            'markitdown_available': True
        })
    except Exception as e:
        return jsonify({
            'tesseract_available': False,
            'error': str(e),
            'markitdown_available': True
        })

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)